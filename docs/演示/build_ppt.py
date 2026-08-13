#!/usr/bin/env python3
# 业务安全态势系统 — 演示 PPT 生成脚本（最新系统：真实数据仪表盘 + 弹出/跳转交互）
# 产出:
#   docs/演示/业务安全态势系统.pptx   综合演示 deck（6 页）
#   docs/演示/业务态势.pptx           业务主题单页
# 设计风格：深蓝底 1A2B4A / 主蓝 1E508C / 绿勾 2E7D32 / 微软雅黑
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DEMO = os.path.normpath(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), '演示'))

BG    = RGBColor(0x1A, 0x2B, 0x4A)
BG2   = RGBColor(0x23, 0x3B, 0x63)
ACC   = RGBColor(0x1E, 0x50, 0x8C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB   = RGBColor(0x9F, 0xB3, 0xCC)
FONT  = 'Microsoft YaHei'
W, H = 10.0, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color=ACC):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def txt(slide, x, y, w, h, runs, size=15, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str): runs = [(runs, color, bold)]
    first = True
    for text, c, b in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c; r.font.name = FONT
    return tb

def bullets(slide, x, y, w, items, size=13, gap=0.36, mark='✓ '):
    y0 = y
    for text in items:
        txt(slide, x, y0, w, 0.35, [(mark, GREEN, True), (text, WHITE, False)], size=size)
        y0 += gap
    return y0

def header(slide, title, subtitle):
    txt(slide, 0.6, 0.22, 8.0, 0.5, title, size=25, bold=True)
    txt(slide, 0.6, 0.68, 8.8, 0.35, subtitle, size=13, color=SUB)
    rect(slide, 0.6, 1.02, 2.2, 0.045)

def shot(slide, name, x, y, w, h=None):
    if h is None: h = w * 9 / 16
    slide.shapes.add_picture(os.path.join(DEMO, f'{name}-1920.png'), Inches(x), Inches(y), Inches(w), Inches(h))

def caption(slide, text, x, y, w, size=10.5):
    txt(slide, x, y, w, 0.3, text, size=size, color=SUB)

# ---------------------------------------------------------------- slide 1 封面
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, W, 0.12, ACC)
txt(s, 0.6, 2.05, 5.0, 1.0, '业务安全态势系统', size=44, bold=True)
txt(s, 0.6, 3.0, 5.4, 0.5, '综合 · 业务 · 运维 · 终端 实时态势演示', size=18, color=SUB)
txt(s, 0.6, 3.5, 5.4, 0.4, '真实数据仪表盘 · 弹出交互 · 专题下钻', size=14, color=SUB)
shot(s, 'overview-main', 5.35, 1.55, 4.1)
rect(s, 0, 6.55, W, 0.75, BG2)
txt(s, 0.6, 6.68, 8.8, 0.5,
    '前端 Vue 3 + TypeScript + Vite + ECharts ｜ 后端 Spring Boot 3 + Java 17 ｜ 数据库 MySQL 8 ｜ nginx 部署',
    size=12, color=SUB)

# ---------------------------------------------------------------- slide 2 综合态势总览
s = prs.slides.add_slide(BLANK); bg(s)
header(s, '综合态势总览', '首页大盘：人员在线、密信收发、区域链路、安全事件与签阅处置一屏聚合')
bullets(s, 0.6, 1.15, 8.8, [
    '在线概览：在线率 65%，密信文件收发 82 份，用户/部门收发排名支持总量、发送、接收切换；',
    '区域链路：北京、阿联酋、新加坡、德国、肯尼亚、巴西六区域链路状态与当日流量（共 9.7 GB）实时呈现；',
    '事件与处置：安全事件滚动更新（高危 2 / 中危 2），签阅处置 54/68、处理率 79.4%、待处理 11 份。'])
shot(s, 'overview-main', 0.6, 2.35, 8.8)

# ---------------------------------------------------------------- slide 3 弹出交互
s = prs.slides.add_slide(BLANK); bg(s)
header(s, '弹出交互', '点击即达：人员详情抽屉、线路安全智能分析提示、线路拓扑面板')
bullets(s, 0.6, 1.15, 8.8, [
    '安全事件条目点击后弹出人员详情抽屉（人员概况 / 装备套件 / 安全事件 / 业务活动四标签）；',
    '线路遭受攻击时弹出"线路安全智能分析"toast，提示分段绕行策略已下发；',
    '点击 toast 打开对应线路拓扑面板，查看各跳节点状态与绕行路径。'])
w3 = 2.9; g3 = 0.15; x0 = 0.6; y3 = 2.4; h3 = w3 * 9 / 16
for i, (img, cap) in enumerate([
    ('overview-drawer-person', '安全事件 → 人员详情抽屉'),
    ('overview-route-toast', '线路安全智能分析 toast'),
    ('overview-topology-panel', '点击 toast → 线路拓扑面板')]):
    x = x0 + i * (w3 + g3)
    shot(s, img, x, y3, w3, h3)
    caption(s, cap, x, y3 + h3 + 0.08, w3)

# ---------------------------------------------------------------- slide 4 业务态势（下钻明细）
s = prs.slides.add_slide(BLANK); bg(s)
header(s, '业务态势 · 专题下钻', '首页卡片一键跳转：密信态势、签阅态势、系统流量三大专题（返回综合态势一键回跳）')
bullets(s, 0.6, 1.15, 8.8, [
    '密信态势：消息/文件收发量、当前登录用户、分时趋势与用户/部门收发排名，支持发送/接收切换；',
    '签阅态势：收到 68 / 已处理 54 / 待处理 11 / 异常退回 3，人员与部门处理排名可视化；',
    '系统流量：业务系统吞吐排名与区域链路流量排名（上行 + 下行），今日流量实时累计。'])
w2 = 4.3; xA = 0.6; xB = 5.1; y4 = 2.4; h4 = w2 * 9 / 16
shot(s, 'business-message', xA, y4, w2, h4)
shot(s, 'business-signing', xB, y4, w2, h4)
caption(s, '密信态势（message）', xA, y4 + h4 + 0.08, w2)
caption(s, '签阅态势（signing）', xB, y4 + h4 + 0.08, w2)
txt(s, 0.6, y4 + h4 + 0.45, 8.8, 0.35, '跳转示意：/overview 点击"今日业务摘要 / 签阅处置"等面板 → /business?topic=message|signing|traffic → 返回综合态势', size=10.5, color=SUB)

# ---------------------------------------------------------------- slide 5 运维态势
s = prs.slides.add_slide(BLANK); bg(s)
header(s, '运维态势', '真实运维域 /ops：Probe / External / Manual 多源接入')
bullets(s, 0.6, 1.15, 8.8, [
    '多源接入：Java Probe 经 /api/ops/ingest 统一接入，支持 external 与 manual 补录；',
    '主机纳管：主机列表、详情与近七日趋势联动，TopN / 白名单进程分析；',
    '告警闭环：最新告警实时上屏，总览卡片与来源概览呈现整体运行态势。'])
shot(s, 'ops-main', 0.6, 2.35, 8.8)

# ---------------------------------------------------------------- slide 6 终端态势
s = prs.slides.add_slide(BLANK); bg(s)
header(s, '终端态势', '终端资产 · 人员关联 · 保障状态')
bullets(s, 0.6, 1.15, 8.8, [
    '终端资产：全球终端按国家聚合展示，规模与状态分布一目了然；',
    '人员关联：终端与责任人员对应，归属确认与保障状态全程可见；',
    '风险处置：高风险终端、异常事件与策略命中情况联动呈现。'])
shot(s, 'terminal-main', 0.6, 2.35, 8.8)

prs.save(os.path.join(DEMO, '业务安全态势系统.pptx'))
print('OK 业务安全态势系统.pptx', len(prs.slides._sldIdLst), 'slides')

# ---------------------------------------------------------------- 业务态势.pptx（单页业务主题）
prs2 = Presentation()
prs2.slide_width = Inches(W); prs2.slide_height = Inches(H)
s = prs2.slides.add_slide(prs2.slide_layouts[6]); bg(s)
rect(s, 0, 0, W, 0.12, ACC)
txt(s, 0.6, 0.3, 8.9, 0.4, '业务态势　·　密信 / 签阅 / 系统流量', size=22, bold=True)
txt(s, 0.6, 0.8, 8.9, 0.3,
    '业务主题态势聚焦密信传输、签阅流转与系统流量三大核心业务，首页卡片可一键下钻并返回。',
    size=13, color=SUB)
bullets(s, 0.6, 1.2, 8.6, [
    '密信态势：消息与文件收发量统计、当前登录用户、分时趋势及用户/部门收发排名实时展示；',
    '签阅态势：收到、已处理、待处理、异常/退回流转状态统计，人员与部门处理排名可视化；',
    '系统流量：业务系统吞吐排名与区域链路流量排名，各系统今日流量实时累计。'],
    size=13, gap=0.45)
rect(s, 0.6, 2.85, 8.9, 0.02, BG2)
shot(s, 'business-message', 1.7, 3.05, 6.6)
txt(s, 0.6, 6.85, 8.9, 0.4,
    '数据为模拟演示场景，实时滚动更新：签阅收到 68 份、处理率 79.4%、待处理 11 份、异常退回 3 份。',
    size=10.5, color=SUB)
prs2.save(os.path.join(DEMO, '业务态势.pptx'))
print('OK 业务态势.pptx')
